# project_management_system/authentication/views.py
import logging
import google.api_core.exceptions as google_exceptions
import requests
import json
import re
import tempfile
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from django.http import HttpResponse
from django.shortcuts import get_object_or_404
import io
from django.utils import timezone
from .models import TypingStatus
from django.core.mail import send_mail # Ensure this is here too
from django.conf import settings
from rest_framework import generics, status, views, serializers
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework.parsers import MultiPartParser, FormParser, JSONParser
from rest_framework.permissions import IsAuthenticated, AllowAny
from project_management.project_analyzer import ProjectAnalyzer
from .models import (
    User, ProjectSubmission, Group, Project, Team, Message, 
    VivaSession, VivaQuestion, ProgressUpdate, ProjectArtifact, 
    Task, CodeReview, TypingStatus, PasswordResetOTP, Checkpoint,
    TimedAssignment, AssignmentSubmission, StudentActivityLog
)
from .permissions import IsTeacherOrAdmin, IsProjectMemberOrTeacher, IsAdminUser
from .serializers import (
    UserSerializer, ProjectSubmissionSerializer, GroupSerializer, 
    ApprovedProjectSerializer, StudentSubmissionSerializer, 
    MessageSerializer, VivaSessionSerializer, VivaQuestionSerializer,
    ProgressUpdateSerializer, PasswordResetRequestSerializer, 
    PasswordResetConfirmSerializer, ProjectArtifactSerializer, 
    TaskSerializer, CodeReviewSerializer, TeacherSubmissionSerializer,
    SimilarProjectSerializer, ProjectSerializer, CheckpointSerializer,
    TimedAssignmentSerializer, AssignmentSubmissionSerializer,
    AlumniProjectSerializer, StudentActivityLogSerializer
)
from django.db.models import Count, Sum, Q


logger = logging.getLogger(__name__)
analyzer = ProjectAnalyzer()

# --- NEW HELPER FUNCTION TO BUILD CONTEXT ---
def _build_project_context(project: Project, user_prompt: str = "") -> str:
    """
    Comprehensive Context Builder.
    Includes:
    1. Project Basics & Student Info
    2. Initial AI Scores
    3. Progress Update History (with Sentiment)
    4. Viva Examination History
    5. Artifacts (Screenshots/Code) - Full text
    6. Documentation (PDF) - AI Feedback + Truncated Text
    """
    # 1. Basic Info
    context = f"PROJECT REPORT (ID: {project.id})\n"
    context += f"Title: {project.title}\n"
    context += f"Status: {project.status}\n"
    context += f"Progress: {project.progress_percentage}%\n"
    context += f"Category: {project.category}\n"
    if project.github_repo_link:
        context += f"GitHub Repo: {project.github_repo_link}\n"
    
    # Include Audit Scores if available
    if project.audit_security_score:
         context += f"Security Score: {project.audit_security_score}/100\n"
    if project.audit_quality_score:
         context += f"Quality Score: {project.audit_quality_score}/100\n"

    # Include Audit Report Summary (Key Findings) if available
    if project.audit_report:
        try:
            # Assuming audit_report is a dict with 'summary' or similar
            summary = project.audit_report.get('summary', '') or project.audit_report.get('executive_summary', '')
            if summary:
                context += f"Latest Audit Summary: {summary}\n"
        except:
            pass

    # 2. Student & Team Info
    student = project.submission.student
    context += f"\nSTUDENT/TEAM DETAILS\n"
    context += f"Lead Student: {student.first_name} {student.last_name} (Username: {student.username})\n"
    
    # Check for Group/Team
    if hasattr(project, 'team'):
        context += f"Team Members:\n"
        for member in project.team.members.all():
            context += f"- {member.first_name} {member.last_name} ({member.username})\n"
    elif hasattr(project.submission, 'group') and project.submission.group:
        context += f"Group: {project.submission.group.name}\n"
        context += f"Group Members (Potential):\n"
        for member in project.submission.group.students.all():
            context += f"- {member.first_name} {member.last_name} ({member.username})\n"
    
    context += f"Email: {student.email}\n"

    # 3. Initial AI Evaluation
    sub = project.submission
    context += f"\nINITIAL PROPOSAL EVALUATION\n"
    context += f"Relevance Score: {sub.relevance_score}/10\n"
    context += f"Feasibility Score: {sub.feasibility_score}/10\n"
    context += f"Innovation Score: {sub.innovation_score}/10\n"
    context += f"Abstract: {sub.abstract_text}\n"

    # 4. PROGRESS UPDATE HISTORY
    # Limit to last 5 updates to prevent Token Overflow (Groq 413)
    progress_logs = ProgressUpdate.objects.filter(project=project).order_by('-created_at')[:5]
    progress_logs = reversed(list(progress_logs)) # Chronological order
    
    context += f"\nPROGRESS UPDATE HISTORY (Last 5 items)\n"

    # We iterate over the list now
    has_logs = False
    for i, log in enumerate(progress_logs, 1):
        has_logs = True
        context += f"\n-- Log {i} ({log.created_at.strftime('%Y-%m-%d')}) --\n"
        # Attribute to specific author
        author_name = log.author.username
        if hasattr(log.author, 'first_name') and log.author.first_name:
            author_name = f"{log.author.first_name} {log.author.last_name}"
        
        context += f"Author: {author_name}\n"
        # Truncate long text
        text = log.update_text if len(log.update_text) < 500 else log.update_text[:500] + "..."
        context += f"Report: {text}\n"
        # Include Sentiment if available
        sentiment = getattr(log, 'sentiment', 'N/A')
        context += f"Sentiment: {sentiment}\n"
        context += f"AI-Suggested Progress: {log.ai_suggested_percentage}%\n"
    
    if not has_logs:
        context += "No progress logs have been submitted yet.\n"

    # 5. VIVA HISTORY
    # Limit to last 15 sessions to cover all team members
    viva_sessions = VivaSession.objects.filter(project=project).order_by('-created_at')[:15]
    viva_sessions = reversed(list(viva_sessions))
    
    context += f"\nVIVA EXAMINATION HISTORY (Last 3 sessions)\n"
    
    has_vivas = False
    for i, session in enumerate(viva_sessions, 1):
        has_vivas = True
        s_name = session.student.username
        if hasattr(session.student, 'first_name') and session.student.first_name:
                s_name = f"{session.student.first_name} {session.student.last_name}"
        
        context += f"\n-- Session {i} ({session.created_at.strftime('%Y-%m-%d')}) - Student: {s_name} --\n"
        total_score = 0
        total_questions = 0
        # Limit details per session
        for q in session.questions.all()[:5]:
            context += f"Q: {q.question_text}\n"
            ans = q.student_answer if q.student_answer else 'Not answered'
            if len(ans) > 200: ans = ans[:200] + "..."
            context += f"A: {ans}\n"
            context += f"Score: {q.ai_score}/10\n"
            context += f"Feedback: {q.ai_feedback}\n"
            if q.ai_score is not None:
                total_score += q.ai_score
                total_questions += 1
        avg_score = (total_score / total_questions) if total_questions > 0 else 0
        context += f"Avg Score for Session {i}: {avg_score:.1f}/10\n"

    if not has_vivas:
        context += "No viva sessions have been attempted yet.\n"

    # 6. PROJECT ARTIFACTS (Screenshots & Code)
    artifacts = ProjectArtifact.objects.filter(project=project).order_by('uploaded_at')
    context += f"\nPROJECT DOCUMENTS & SCREENSHOTS ({artifacts.count()} files)\n"
    if not artifacts.exists():
        context += "No screenshots or documents have been uploaded.\n"
    else:
        for i, art in enumerate(artifacts, 1):
            context += f"\n-- Document {i} (Uploaded: {art.uploaded_at.strftime('%Y-%m-%d')}) --\n"
            context += f"Description: {art.description}\n"
            context += f"AI Auto-Tags: {art.ai_tags}\n"
            context += f"EXTRACTED TEXT CONTENT:\n{art.extracted_text}\n"

    # 7. DOCUMENTATION (PDF Report)
    if project.final_report:
        context += f"\nFULL PROJECT DOCUMENTATION (PDF)\n"
        context += f"Status: Uploaded\n"
        
        # A. The AI Critique
        if project.ai_report_feedback:
             context += f"AI CRITIQUE:\n{project.ai_report_feedback}\n"
        
        # B. The Actual Content (Truncated to 30k chars to prevent errors)
        if project.final_report_content:
             context += f"PDF FULL TEXT CONTENT (First 30k chars):\n"
             context += f"{project.final_report_content[:30000]}\n"
             if len(project.final_report_content) > 30000:
                context += "\n...(Text truncated for API limits)...\n"
    else:
        context += "\nNo full documentation uploaded.\n"

    # 8. CODE REVIEWS (NEW)
    reviews = CodeReview.objects.filter(project=project).order_by('-uploaded_at')[:3]
    if reviews.exists():
        context += "\n--- RECENT CODE REVIEWS ---\n"
        for rev in reviews:
            context += f"- File: {rev.file_name} (Score: Sec {rev.security_score}/10, Qual {rev.quality_score}/10)\n"
            context += f"  AI Feedback: {rev.ai_feedback[:200]}...\n"

    return context

class ProjectSubmissionView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser, JSONParser,)

    def post(self, request, *args, **kwargs):
        user = request.user
        if user.is_anonymous:
             return Response({"error": "User must be logged in."}, status=status.HTTP_401_UNAUTHORIZED)
        
        force_submit = request.data.get('force_submit') == 'true' 

        # 1. Prepare the data that needs validation
        # We only validate what the user provides
        validation_data = {
            'title': request.data.get('title', '').strip(),
            'abstract_text': request.data.get('abstract_text', '').strip(),
            'abstract_file': request.FILES.get('abstract_file'),
            'audio_file': request.FILES.get('audio_file'),
        }

        # Find and add the group
        student_groups = list(user.student_groups.all())
        if not student_groups:
            return Response({"error": "You must be a member of a group to submit a project."}, status=status.HTTP_400_BAD_REQUEST)
        validation_data['group'] = student_groups[0].id
        
        # This text is used for all AI analysis
        text_to_analyze = validation_data['title'] if not validation_data['abstract_text'] else validation_data['abstract_text']

        # 2. Prepare the data that will be SAVED (not validated)
        # These are all the AI-generated fields
        save_kwargs = {}

        if not force_submit:
            # --- STAGE 1: Analysis Path ---
            if not text_to_analyze:
                return Response({"error": "Abstract text or title is required."}, status=status.HTTP_400_BAD_REQUEST)

            try:
                existing_submissions = ProjectSubmission.objects.all().values('abstract_text', 'title', 'student__username', 'logical_fingerprint')
                
                analysis_result = analyzer.check_plagiarism_and_suggest_features(
                    title=validation_data['title'], 
                    abstract=text_to_analyze,
                    existing_submissions=list(existing_submissions)
                )
            except google_exceptions.ResourceExhausted as e:
                return Response(
                    {"detail": "The AI analyzer is temporarily busy. Please try again in one minute."},
                    status=status.HTTP_429_TOO_MANY_REQUESTS
                )
            except Exception as e:
                return Response(
                    {"detail": "An unexpected error occurred during AI analysis."},
                    status=status.HTTP_500_INTERNAL_SERVER_ERROR
                )

            # Check for plagiarism
            if analysis_result['originality_status'] == "BLOCKED_HIGH_SIMILARITY":
                similar_project_data = {}
                if analysis_result['most_similar_project']:
                    analysis_result['most_similar_project']['student__username'] = analysis_result['most_similar_project'].get('student', 'N/A')
                    similar_project_data = SimilarProjectSerializer(analysis_result['most_similar_project']).data

                return Response({
                    "detail": "Submission Blocked: High Semantic Similarity Detected.",
                    "suggestions": analysis_result.get('suggested_features'),
                    "similar_project": similar_project_data,
                    'relevance_score': analysis_result.get('relevance', 0),
                    'feasibility_score': analysis_result.get('feasibility', 0),
                    'innovation_score': analysis_result.get('innovation', 0)
                }, status=status.HTTP_409_CONFLICT)
            
            # If OK, add AI results to our save_kwargs
            save_kwargs['relevance_score'] = analysis_result['relevance']
            save_kwargs['feasibility_score'] = analysis_result['feasibility']
            save_kwargs['innovation_score'] = analysis_result['innovation']
            save_kwargs['ai_summary'] = analysis_result.get('full_report') # Explicitly save the report as summary
            save_kwargs['ai_similarity_report'] = analysis_result.get('most_similar_project')
            save_kwargs['ai_suggested_features'] = analysis_result.get('suggested_features')
            save_kwargs['logical_fingerprint'] = analysis_result.get('logical_fingerprint')
        
        else:
            # --- STAGE 3: Force Submit Path ---
            # Get the AI data back from the request
            save_kwargs['relevance_score'] = request.data.get('relevance_score', 0)
            save_kwargs['feasibility_score'] = request.data.get('feasibility_score', 0)
            save_kwargs['innovation_score'] = request.data.get('innovation_score', 0)
            save_kwargs['ai_suggested_features'] = request.data.get('ai_suggested_features')
            try:
                fp_str = request.data.get('logical_fingerprint')
                if fp_str and fp_str != 'null':
                     save_kwargs['logical_fingerprint'] = json.loads(fp_str) if isinstance(fp_str, str) else fp_str
            except:
                pass
            
            try:
                report_str = request.data.get('ai_similarity_report')
                if report_str and report_str != 'null':
                    save_kwargs['ai_similarity_report'] = json.loads(report_str)
            except json.JSONDecodeError:
                save_kwargs['ai_similarity_report'] = None

        # --- Common Save Logic ---
        
        # 3. Validate ONLY the user's data
        serializer = ProjectSubmissionSerializer(data=validation_data)
        if not serializer.is_valid():
            # This is the 400 error.
            return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

        # 4. Call microservices (this is where your logs were correct)
        print("--- Calling AI Microservice ---") # Added for your debugging
        try:
            response_tags = requests.post("http://127.0.0.1:8001/extract-keywords", json={"text": text_to_analyze}, timeout=5)
            if response_tags.status_code == 200: 
                save_kwargs['tags'] = response_tags.json().get('keywords')
        except Exception as e:
            print(f"Keyword AI call failed: {e}")
        
        try:
            response_summary = requests.post("http://127.0.0.1:8001/summarize", json={"text": text_to_analyze}, timeout=60)
            if response_summary.status_code == 200: 
                save_kwargs['ai_summary'] = response_summary.json().get('summary')
        except Exception as e:
            print(f"An unknown error occurred during summary AI call: {e}")
        
        # 5. Save the project
        new_embedding = analyzer.get_embedding(text_to_analyze)
        
        serializer.save(
            student=user, 
            embedding=new_embedding,
            transcribed_text=None,
            **save_kwargs  # Pass all AI data here
        )
        
        return Response(serializer.data, status=status.HTTP_201_CREATED)


class TeacherDashboardView(APIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]

    def get(self, request, submission_id=None, *args, **kwargs):
        if submission_id:
            try:
                submission = ProjectSubmission.objects.get(id=submission_id)
            except ProjectSubmission.DoesNotExist:
                return Response({"detail": "Submission not found."}, status=status.HTTP_404_NOT_FOUND)
            
            # Since teachers can view "Other Ongoing Projects", we allow viewing details of any submission
            # provided they are authorized as a teacher.
            serializer = TeacherSubmissionSerializer(submission)
            return Response(serializer.data, status=status.HTTP_200_OK)

        user = request.user
        # Strict Teacher Filter: Only show submissions from assigned teaching groups.
        teacher_groups = user.teaching_groups.all()
        queryset = ProjectSubmission.objects.filter(group__in=teacher_groups).order_by('-submitted_at')

        # We want to show ALL submissions, even if the group has a project.
        # This allows teachers to see if a group is trying to pivot or switch projects.
        # However, checking "Single Project Per Group" at APPROVAL time is still good practice.
        
        # groups_with_projects = Project.objects.filter(
        #     status__in=['In Progress', 'Completed'],
        #     submission__group__isnull=False
        # ).values_list('submission__group_id', flat=True)
        
        # submissions = queryset.filter(status='Submitted').exclude(group_id__in=groups_with_projects)
        submissions = queryset.filter(status='Submitted')
            
        serializer = TeacherSubmissionSerializer(submissions, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    def patch(self, request, submission_id, *args, **kwargs):
        try:
            submission = ProjectSubmission.objects.get(id=submission_id)
        except ProjectSubmission.DoesNotExist:
            return Response({"detail": "Submission not found."}, status=status.HTTP_404_NOT_FOUND)
        teacher_groups = request.user.teaching_groups.all()
        if submission.group not in teacher_groups:
            return Response({"detail": "You do not have permission to review this project."}, status=status.HTTP_403_FORBIDDEN)
        new_status = request.data.get('status')
        if new_status not in ['Approved', 'Rejected']:
            return Response({"detail": "Invalid status. Must be 'Approved' or 'Rejected'."}, status=status.HTTP_400_BAD_REQUEST)
        if submission.status != 'Submitted':
            return Response({"detail": "This project has already been reviewed."}, status=status.HTTP_400_BAD_REQUEST)
        submission.status = new_status
        submission.save()
        if new_status == 'Approved':
            # Check if group already has an active project
            if submission.group:
                active_projects = Project.objects.filter(
                    submission__group=submission.group,
                    status__in=['In Progress', 'Completed']
                )
                if active_projects.exists():
                    return Response(
                        {"detail": "This group already has an active or completed project. Only one project is allowed per group."},
                        status=status.HTTP_400_BAD_REQUEST
                    )

            project = Project.objects.create(
                submission=submission, title=submission.title,
                abstract=submission.abstract_text, status='In Progress'
            )
            team = Team.objects.create(project=project)
            # Add ALL group members to the team
            if submission.group:
                for student in submission.group.students.all():
                    team.members.add(student)
            else:
                # Fallback for individual submissions (if allowed)
                team.members.add(submission.student)
        serializer = TeacherSubmissionSerializer(submission)
        return Response(serializer.data, status=status.HTTP_200_OK)


class StudentDashboardView(APIView):
    permission_classes = [IsAuthenticated]
    def get(self, request, *args, **kwargs):
        user = request.user
        student_groups = user.student_groups.all()
        # Fetch submissions if:
        # 1. User is in the submission's initial group
        # 2. User is the submitter
        # 3. User is a member of the project's team (for Approved/In Progress projects)
        submissions = ProjectSubmission.objects.filter(
            Q(group__in=student_groups) | 
            Q(student=user) | 
            Q(project__team__members=user)
        ).distinct().order_by('-submitted_at')
        serializer = StudentSubmissionSerializer(submissions, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)


class AIChatbotView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser, JSONParser,)

    def post(self, request, *args, **kwargs):
        user_prompt = request.data.get('prompt')
        audio_file = request.data.get('audio_file')
        if not user_prompt and not audio_file:
            return Response({"error": "Prompt or audio file not provided."}, status=status.HTTP_400_BAD_REQUEST)
        if audio_file:
            # audio_file assumed to be an UploadedFile; transcribe similarly to submission flow
            try:
                if hasattr(audio_file, "temporary_file_path"):
                    user_prompt = analyzer.transcribe_audio(audio_file.temporary_file_path())
                else:
                    with tempfile.NamedTemporaryFile(delete=True, suffix=".wav") as tf:
                        for chunk in audio_file.chunks():
                            tf.write(chunk)
                        tf.flush()
                        user_prompt = analyzer.transcribe_audio(tf.name)
            except Exception as e:
                logger.error(f"Audio transcription in chatbot failed: {e}")
                return Response({"error": "Failed to transcribe audio."}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
            if not user_prompt:
                return Response({"error": "Failed to transcribe audio."}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        # Context Building
        project_id = request.data.get('project_id')
        project_context = ""

        try:
            project_submission = None
            
            if project_id:
                try:
                    project_submission = ProjectSubmission.objects.get(id=project_id)
                except ProjectSubmission.DoesNotExist:
                     pass
            
            if not project_submission:
                # 1. Prioritize Active/Completed Projects where user is a team member
                active_project = Project.objects.filter(team__members=request.user).order_by('-submission__submitted_at').first()
                if active_project:
                    project_submission = active_project.submission
                else:
                    # 1b. Fallback: Check if user is in a Group that has a project (even if Team object is missing)
                    group_project = Project.objects.filter(submission__group__students=request.user).order_by('-submission__submitted_at').first()
                    if group_project:
                         project_submission = group_project.submission
                    else:
                        # 2. Final Fallback: Get latest submission by this user (e.g. pending proposal/no group yet)
                        project_submission = ProjectSubmission.objects.filter(student=request.user).order_by('-submitted_at').first()
            
            if project_submission:
                # Determine role
                role = "Team Member"
                if project_submission.student == request.user:
                    role = "Project Leader"
                    
                status = project_submission.status
                project_context = f"Project Title: {project_submission.title}\n"
                project_context += f"Status: {status}\n"
                project_context += f"User Role: {role}\n"
                project_context += f"Abstract: {project_submission.abstract_text}"
        except Exception as e:
            logger.warning(f"Could not load project context for chat: {e}")

        # Call AI with context
        ai_response = analyzer.get_chat_response(user_prompt, context=project_context)
        return Response({"response": ai_response}, status=status.HTTP_200_OK)


class AIVivaView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, *args, **kwargs):

        project_id = request.data.get('project_id')
        if not project_id:
            return Response({"error": "Project ID required."}, status=status.HTTP_400_BAD_REQUEST)


        project = get_object_or_404(Project, id=project_id)

        # Call AI Microservice (MCP Endpoint)
        try:
             MICROSERVICE_URL = "http://127.0.0.1:8001/mcp-viva-questions"
             # Resolution of student: user is request.user (Student)
             payload = {
                 "student_username": request.user.username
             }
             
             response = requests.post(MICROSERVICE_URL, json=payload, timeout=45)
             
             if response.status_code != 200:
                  return Response({"error": f"AI Error: {response.text}"}, status=response.status_code)
                  
             data = response.json()
             if "questions" in data:
                 questions_text_list = data["questions"]
             elif "response" in data:
                  # Fallback if it returns raw text, try to parse or treat as single question?
                  # Ideally endpoint returns {"questions": ["Q1", "Q2"]}
                  # Let's assume it does. If not, use legacy fallback or error.
                  questions_text_list = [data["response"]]
             else:
                  questions_text_list = ["Could not generate specific questions. Please describe your project status."]

        except Exception as e:
             return Response({"error": f"AI Connection Failed: {str(e)}"}, status=500)

        session = VivaSession.objects.create(project=project, student=request.user)
        viva_questions = []
        for q_text in questions_text_list:
            viva_questions.append(VivaQuestion(session=session, question_text=q_text))
        VivaQuestion.objects.bulk_create(viva_questions)

        serializer = VivaSessionSerializer(session)
        return Response(serializer.data, status=status.HTTP_201_CREATED)


class AIVivaEvaluationView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, *args, **kwargs):
        question_id = request.data.get('question_id')
        student_answer = request.data.get('answer')
        if not question_id or not student_answer:
            return Response({"error": "Question ID and answer are required."}, status=status.HTTP_400_BAD_REQUEST)
        viva_question = get_object_or_404(VivaQuestion, id=question_id)
        if viva_question.session.student != request.user:
            return Response({"error": "Unauthorized."}, status=status.HTTP_403_FORBIDDEN)
        abstract = viva_question.session.project.abstract
        evaluation_result = analyzer.evaluate_viva_answer(
            question=viva_question.question_text,
            answer=student_answer,
            abstract=abstract
        )
        try:
            score_int = int(float(evaluation_result.get('score', 0)))
        except (ValueError, TypeError):
            score_int = 0
        viva_question.student_answer = student_answer
        viva_question.ai_score = score_int
        viva_question.ai_feedback = evaluation_result.get('feedback', '')
        viva_question.save()
        return Response(VivaQuestionSerializer(viva_question).data, status=status.HTTP_200_OK)


class ProjectArchiveView(APIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]

    def patch(self, request, project_id, *args, **kwargs):
        try:
            project = Project.objects.get(id=project_id)
        except Project.DoesNotExist:
            return Response({"detail": "Project not found."}, status=status.HTTP_404_NOT_FOUND)
        new_status = request.data.get('status')
        if new_status not in ['Completed', 'Archived']:
            return Response({"detail": "Invalid status. Must be 'Completed' or 'Archived'."}, status=status.HTTP_400_BAD_REQUEST)
        if new_status == 'Completed':
            if project.status == 'In Progress':
                project.status = 'Completed'
                project.save()
            else:
                return Response({"detail": "Project must be 'In Progress' to be marked as 'Completed'."}, status=status.HTTP_400_BAD_REQUEST)
        elif new_status == 'Archived':
            if project.status == 'Completed':
                project.status = 'Archived'
                project.save()
            else:
                return Response({"detail": "Project must be 'Completed' to be archived."}, status=status.HTTP_400_BAD_REQUEST)
        submission = project.submission
        submission.status = new_status
        submission.save()

        # Handle Alumni Conversion
        if new_status == 'Completed':
            project.is_alumni = True
            
            # Trigger AI Analysis for Alumni Ranking logic
            try:
                # We can reuse extraction logic or build a new one
                # For now, let's just generate random or basic scores if AI service isn't specialized yet
                # ideally: analyzer.analyze_alumni_impact(project)
                
                # Mocking the AI call for trend/relevance for now, or use existing scores
                if hasattr(submission, 'innovation_score'):
                     # normalize or adjust
                     pass

                # Let's say we want to refresh these scores based on "Current Market Trends"
                # For this MVP, we will rely on the initial submission scores + Audit scores
                # But we set the flag so they show up in the Alumni list.
                pass
            except Exception as e:
                print(f"Alumni conversion error: {e}")
            
            project.save()

        return Response({"detail": f"Project status updated to {new_status}."}, status=status.HTTP_200_OK)





class AnalyticsView(generics.ListAPIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]

    def get_queryset(self):
        return None

    def list(self, request, *args, **kwargs):
        status_counts = Project.objects.values('status').annotate(count=Count('status'))
        category_counts = Project.objects.values('category').annotate(count=Count('category'))
        top_innovative = Project.objects.filter(status='Completed').order_by('-submission__innovation_score')[:5]
        top_innovative_data = [{'title': p.title, 'score': p.submission.innovation_score} for p in top_innovative]
        data = {
            'project_status_counts': list(status_counts),
            'project_category_counts': list(category_counts),
            'top_innovative_projects': top_innovative_data,
        }
        return Response(data, status=status.HTTP_200_OK)


class LeaderboardView(generics.ListAPIView):
    permission_classes = [IsAuthenticated]
    serializer_class = UserSerializer

    def get_queryset(self):
        queryset = User.objects.annotate(
            total_innovation=Sum('active_projects__project__submission__innovation_score')
        )
        return queryset.filter(
            total_innovation__isnull=False,
            active_projects__project__status='Completed'
        ).order_by('-total_innovation')[:10]


class AlumniPortalView(generics.ListAPIView):
    permission_classes = [IsAuthenticated]
    serializer_class = AlumniProjectSerializer

    def get_queryset(self):
        user = self.request.user
        return Project.objects.filter(
            Q(submission__student=user) | Q(team__members=user),
            status__in=['Completed', 'Archived']
        ).distinct().order_by('-submission__submitted_at')


import numpy as np

class AlumniProjectSearchView(APIView):
    permission_classes = [AllowAny]

    def get(self, request, *args, **kwargs):
        query = request.query_params.get('q', '').strip()
        if not query:
            return Response([], status=status.HTTP_200_OK)

        # Generate embedding for query
        try:
            query_embedding = analyzer.get_embedding(query)
            if not query_embedding:
                 return Response({"error": "AI service unavailable"}, status=status.HTTP_503_SERVICE_UNAVAILABLE)
        except Exception:
             return Response({"error": "AI service error"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        # Fetch candidates
        projects = Project.objects.filter(status__in=['Completed', 'Archived']).select_related('submission')
        
        results = []
        
        def cosine_similarity(v1, v2):
            return np.dot(v1, v2) / (np.linalg.norm(v1) * np.linalg.norm(v2))

        # 1. Keyword Search (Title Match) - Give high score (1.0)
        keyword_matches = Project.objects.filter(
            title__icontains=query, 
            status__in=['Completed', 'Archived']
        )
        seen_ids = set()
        
        for p in keyword_matches:
            results.append((1.0, p))
            seen_ids.add(p.id)

        # 2. Semantic Search
        for p in projects:
            if p.id in seen_ids:
                continue # Skip if already found by keyword
                
            if not p.submission or not p.submission.embedding:
                continue
            
            try:
                emb = p.submission.embedding
                score = cosine_similarity(query_embedding, emb)
                
                # Threshold for "relevance"
                if score > 0.45: 
                    results.append((score, p))
            except Exception:
                continue
        
        # Sort by similarity score descending
        results.sort(key=lambda x: x[0], reverse=True)
        
        # Return top 20 matches
        top_projects = [x[1] for x in results[:20]]
        serializer = AlumniProjectSerializer(top_projects, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)


class AllProjectsView(generics.ListAPIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]
    queryset = Project.objects.all()
    serializer_class = ProjectSerializer


class AdminDashboardView(APIView):
    permission_classes = [IsAuthenticated, IsAdminUser]

    def get(self, request, *args, **kwargs):
        users = User.objects.all()
        groups = Group.objects.all()
        user_serializer = UserSerializer(users, many=True)
        group_serializer = GroupSerializer(groups, many=True)
        return Response({'users': user_serializer.data, 'groups': group_serializer.data}, status=status.HTTP_200_OK)

    def patch(self, request, group_id, *args, **kwargs):
        return Response({"detail": "Group update not implemented yet."}, status=status.HTTP_200_OK)

    def post(self, request, *args, **kwargs):
        """
        Create a new Group.
        # Fixed typo: Checks if name is provided
        """
        name = request.data.get('name')
        description = request.data.get('description', '')

        if not name:
            return Response({"error": "Group name is required."}, status=status.HTTP_400_BAD_REQUEST)
        
        if Group.objects.filter(name=name).exists():
            return Response({"error": "Group with this name already exists."}, status=status.HTTP_400_BAD_REQUEST)

        group = Group.objects.create(name=name, description=description)
        return Response(GroupSerializer(group).data, status=status.HTTP_201_CREATED)


class AppointedTeacherDashboard(generics.ListAPIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]
    serializer_class = TeacherSubmissionSerializer

    def get_queryset(self):
        # 1. Strict Teacher Filter: Only show submitted projects for assigned groups.
        teacher_groups = self.request.user.teaching_groups.all()
        queryset = ProjectSubmission.objects.filter(
            group__in=teacher_groups,
            status='Submitted'
        )

        # 2. Exclude submissions for groups that already have an 'In Progress' or 'Completed' project
        # 2. Exclude submissions for groups that already have an 'In Progress' or 'Completed' project
        # UPDATE: We now allow seeing them so teachers can manage the transition.
        # groups_with_projects = Project.objects.filter(
        #     status__in=['In Progress', 'Completed'],
        #     submission__group__isnull=False
        # ).values_list('submission__group_id', flat=True)

        # return queryset.exclude(group_id__in=groups_with_projects).order_by('-submitted_at')
        return queryset.order_by('-submitted_at')


class UnappointedTeacherDashboard(generics.ListAPIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]
    serializer_class = TeacherSubmissionSerializer

    def get_queryset(self):
        # 1. Start with projects that have no teacher
        queryset = ProjectSubmission.objects.filter(
            group__teachers__isnull=True,
            status='Submitted'
        )
        
        # 2. Exclude submissions for groups that already have an 'In Progress' or 'Completed' project
        # UPDATE: We now allow seeing them so teachers can manage the transition.
        # groups_with_projects = Project.objects.filter(
        #     status__in=['In Progress', 'Completed'],
        #     submission__group__isnull=False
        # ).values_list('submission__group_id', flat=True)
        
        # return queryset.exclude(group_id__in=groups_with_projects).order_by('-submitted_at')
        return queryset.order_by('-submitted_at')


class ProjectProgressView(views.APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request, project_id, *args, **kwargs):
        try:
            project = Project.objects.get(submission__id=project_id)
            return Response({"progress_percentage": project.progress_percentage, "project_id": project.id}, status=status.HTTP_200_OK)
        except Project.DoesNotExist:
            return Response({"progress_percentage": 0, "project_id": None}, status=status.HTTP_200_OK)


class ProjectLogUpdateView(APIView):
    """
    Handles a student POSTing a new text progress update.
    The AI analyzes the text, determines a new percentage,
    saves the log, and updates the main Project's progress.
    """
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id, *args, **kwargs):
        project = get_object_or_404(Project, id=project_id)
        
        # --- PERMISSION FIX ---
        # Allow if user is the submitter OR if the user is in the project's group
        is_submitter = project.submission.student == request.user
        is_group_member = False
        if project.submission.group:
            is_group_member = project.submission.group.students.filter(id=request.user.id).exists()
            
        # Check if user is in the official project team
        is_team_member = project.team.members.filter(id=request.user.id).exists()

        if not (is_submitter or is_group_member or is_team_member):
            return Response({"error": "You do not have permission to update this project."}, status=status.HTTP_403_FORBIDDEN)
        # ----------------------
        
        update_text = request.data.get('update_text')
        if not update_text:
            return Response({"error": "update_text is required."}, status=status.HTTP_400_BAD_REQUEST)

        # 1. Call for AI-suggested percentage
        ai_percentage = analyzer.analyze_progress_update(
            project_abstract=project.abstract,
            update_text=update_text
        )

        # 2. Call for Sentiment Analysis
        sentiment = None
        try:
            response = requests.post(
                "http://127.0.0.1:8001/sentiment", 
                json={"text": update_text},
                timeout=10 
            )
            if response.status_code == 200:
                sentiment = response.json().get('sentiment')
            else:
                print(f"AI microservice (sentiment) error: {response.text}")
                
        except requests.ConnectionError:
            print("Error: Could not connect to AI microservice at port 8001.")
        except Exception as e:
            print(f"An unknown error occurred during sentiment AI call: {e}")


        # 3. Save the new log entry
        log_entry = ProgressUpdate.objects.create(
            project=project,
            author=request.user,
            update_text=update_text,
            ai_suggested_percentage=ai_percentage,
            sentiment=sentiment
        )

        # 4. Update the main Project's percentage
        project.progress_percentage = ai_percentage
        project.save()

        # Return the newly created log entry
        serializer = ProgressUpdateSerializer(log_entry)
        return Response(serializer.data, status=status.HTTP_201_CREATED)

class TopAlumniProjectsView(generics.ListAPIView):
    permission_classes = [AllowAny]
    serializer_class = AlumniProjectSerializer

    def get_queryset(self):
        # We now filter from Project model where is_alumni=True (or status=Completed)
        # Prioritizing: Innovation > Relevance > Current Trend (Recency)
        return Project.objects.filter(
            status__in=['Completed', 'Archived']
        ).order_by('-submission__innovation_score', '-submission__relevance_score', '-submission__submitted_at')


class ApprovedProjectsView(generics.ListAPIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]
    serializer_class = ApprovedProjectSerializer

    def get_queryset(self):
        user = self.request.user
        
        # Base filter for status
        projects = Project.objects.filter(
            status__in=['In Progress', 'Completed', 'Archived']
        )

        # Strict Teacher Filter: Only show projects from their assigned groups
        teacher_groups = user.teaching_groups.all()
        return projects.filter(submission__group__in=teacher_groups).order_by('-submission__submitted_at')

class ProjectMessagesView(generics.ListCreateAPIView):
    serializer_class = MessageSerializer
    permission_classes = [IsAuthenticated]

    def get_queryset(self):
        project_id = self.kwargs.get('project_id')
        user = self.request.user
        
        # Filter Params
        chat_type = self.request.query_params.get('type') 
        target_user_id = self.request.query_params.get('target_user_id')

        # Base Query: Messages for this project
        queryset = Message.objects.filter(project_id=project_id)

        if chat_type == 'DM' and target_user_id:
            # DM Logic: Fetch messages between Me and Target
            # (We explicitly check Sender OR Recipient here because DMs don't have self-copies)
            queryset = queryset.filter(
                message_type='DM'
            ).filter(
                Q(sender=user, recipient_id=target_user_id) | 
                Q(sender_id=target_user_id, recipient=user)
            )
        
        elif chat_type == 'TEAM_GROUP':
            # Team Logic: Fetch only messages meant for me in this channel
            # This relies on the "Self Copy" we created in perform_create.
            # We DO NOT check 'sender=user' here, to avoid fetching copies sent to others.
            queryset = queryset.filter(
                message_type='TEAM_GROUP',
                recipient=user 
            )

        else:
            # Guide Group Logic (Default)
            # Same as above: Only fetch messages where I am the recipient (includes my self-copy)
            queryset = queryset.filter(
                message_type='GUIDE_GROUP',
                recipient=user
            )

        return queryset.order_by('timestamp')

    def perform_create(self, serializer):
        project_id = self.kwargs.get('project_id')
        project = get_object_or_404(Project, id=project_id)
        sender = self.request.user
        
        message_type = self.request.data.get('message_type', 'GUIDE_GROUP')
        content = serializer.validated_data['content']
        
        created_msgs = []

        if message_type == 'DM':
            # 1. Direct Message
            recipient_id = self.request.data.get('recipient_id')
            if not recipient_id:
                raise serializers.ValidationError("Recipient required for DM.")
            
            recipient = get_object_or_404(User, id=recipient_id)
            msg = Message.objects.create(
                project=project, sender=sender, recipient=recipient,
                content=content, message_type='DM'
            )
            created_msgs.append(msg)

        else:
            # 2. Group Broadcast (GUIDE_GROUP or TEAM_GROUP)
            recipients = []
            
            # Get all potential members
            all_members = []
            if project.submission and project.submission.group:
                all_members = list(project.submission.group.students.all()) + list(project.submission.group.teachers.all())
            elif hasattr(project, 'team'):
                all_members = list(project.team.members.all())

            # Filter based on type
            for member in all_members:
                if member.id == sender.id: continue # Don't send to self

                if message_type == 'TEAM_GROUP':
                    # Only send to Students
                    if member.role == 'Student':
                        recipients.append(member)
                else: 
                    # GUIDE_GROUP -> Send to everyone
                    recipients.append(member)

            # Create copies
            for r in recipients:
                msg = Message.objects.create(
                    project=project, sender=sender, recipient=r,
                    content=content, message_type=message_type
                )
                created_msgs.append(msg)
            
            # Create a copy for SELF so it shows up in my query
            # (Self-copy is useful for the 'recipient=user' filter to work for my own sent msgs)
            self_copy = Message.objects.create(
                project=project, sender=sender, recipient=sender, # Self-message
                content=content, message_type=message_type, is_read=True
            )
            created_msgs.append(self_copy)

        self._created_messages = created_msgs

    def create(self, request, *args, **kwargs):
        # Standard create wrapper
        serializer = self.get_serializer(data=request.data, context={'request': request})
        serializer.is_valid(raise_exception=True)
        self.perform_create(serializer)
        
        # Return just one for the response
        if hasattr(self, '_created_messages') and self._created_messages:
             return Response(self.get_serializer(self._created_messages[0]).data, status=status.HTTP_201_CREATED)
        return Response({}, status=status.HTTP_201_CREATED)
        
class MarkMessagesReadView(APIView):
    """
    Marks all messages in a project as read for the current user.
    """
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        # Filter messages where I am the recipient and they are unread
        unread_msgs = Message.objects.filter(
            project=project, 
            recipient=request.user, 
            is_read=False
        )
        count = unread_msgs.count()
        unread_msgs.update(is_read=True)
        
        return Response({"detail": f"Marked {count} messages as read."}, status=status.HTTP_200_OK)
    
class ProjectVivaListView(generics.ListAPIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]
    serializer_class = VivaSessionSerializer

    def get_queryset(self):
        project_id = self.kwargs.get('project_id')
        project = get_object_or_404(Project, id=project_id)
        return VivaSession.objects.filter(project=project).order_by('-created_at')


# --- NEW VIEW FOR CONTEXT-AWARE CHAT ---
# In authentication/views.py

class ProjectInquiryView(APIView):
    """
    Allows teachers to ask AI questions about a specific project.
    """
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]

    def post(self, request, *args, **kwargs):
        project_id = request.data.get('project_id')
        user_prompt = request.data.get('prompt')

        if not project_id or not user_prompt:
             return Response({"error": "Project ID and prompt are required."}, status=status.HTTP_400_BAD_REQUEST)

        project = get_object_or_404(Project, id=project_id)
        
        # Resolve Student Username from Project
        # Priorities: Team Member -> Submission Student -> Unknown
        student_username = "Unknown"
        if hasattr(project, 'submission') and project.submission and project.submission.student:
            student_username = project.submission.student.username
        elif hasattr(project, 'team') and project.team and project.team.members.exists():
            student_username = project.team.members.first().username
            
        if student_username == "Unknown":
            return Response({"error": "Could not identify a student for this project context."}, status=400)

        # Call AI Microservice (MCP Endpoint)
        try:
            payload = {
                "student_username": student_username,
                "user_message": user_prompt
            }
            # Note: Port 8001 is the AI Microservice
            response = requests.post("http://127.0.0.1:8001/mcp-teacher-chat", json=payload, timeout=60)
            
            if response.status_code == 200:
                data = response.json()
                return Response(data, status=status.HTTP_200_OK)
            else:
                return Response({"error": f"AI Error: {response.text}"}, status=response.status_code)

        except Exception as e:
            return Response({"error": f"AI Microservice failed: {str(e)}"}, status=500)


class AdminUserRoleView(APIView):
    """
    Allows HOD/Admins to update the role of any user.
    """
    permission_classes = [IsAuthenticated, IsAdminUser]

    def patch(self, request, user_id, *args, **kwargs):
        user_to_update = get_object_or_404(User, id=user_id)
        new_role = request.data.get('role')

        if not new_role:
            return Response({"error": "New 'role' is required."}, status=status.HTTP_400_BAD_REQUEST)

        valid_roles = [choice[0] for choice in User.ROLE_CHOICES]
        if new_role not in valid_roles:
            return Response({"error": f"Invalid role. Must be one of {valid_roles}."}, status=status.HTTP_400_BAD_REQUEST)

        user_to_update.role = new_role
        user_to_update.save()

        serializer = UserSerializer(user_to_update)
        return Response(serializer.data, status=status.HTTP_200_OK)


class AdminGroupManagementView(APIView):
    """
    Allows HOD/Admins to add or remove students/teachers from a group.
    """
    permission_classes = [IsAuthenticated, IsAdminUser]

    def patch(self, request, group_id, *args, **kwargs):
        group = get_object_or_404(Group, id=group_id)

        user_id = request.data.get('user_id')
        action = request.data.get('action')  # e.g., 'add_student', 'remove_student', 'add_teacher', 'remove_teacher'

        if not user_id or not action:
            return Response({"error": "'user_id' and 'action' are required."}, status=status.HTTP_400_BAD_REQUEST)

        user = get_object_or_404(User, id=user_id)

        try:
            if action == 'add_student':
                group.students.add(user)
                # Attempt to find if this group has active projects and add student to the Team
                # This depends on your business logic. For now, we will add to ANY team associated with this group
                # OR, more safely, check if the student has a submission in this group that became a project
                from .models import ProjectSubmission, Project, Team
                # Try to find a project submission for this student in this group
                submissions = ProjectSubmission.objects.filter(group=group, student=user)
                for sub in submissions:
                    if hasattr(sub, 'project') and hasattr(sub.project, 'team'):
                         sub.project.team.members.add(user)
                
                # Check for other projects in this group (Shared Group Projects logic)
                # If the group represents a single team, we might want to add them to ALL projects of this group
                # For now, let's look for any project submission linked to this group
                group_submissions = ProjectSubmission.objects.filter(group=group)
                for g_sub in group_submissions:
                     if hasattr(g_sub, 'project') and hasattr(g_sub.project, 'team'):
                          g_sub.project.team.members.add(user)


            elif action == 'remove_student':
                group.students.remove(user)
                # Remove from teams as well
                from .models import ProjectSubmission
                group_submissions = ProjectSubmission.objects.filter(group=group)
                for g_sub in group_submissions:
                     if hasattr(g_sub, 'project') and hasattr(g_sub.project, 'team'):
                          g_sub.project.team.members.remove(user)

            elif action == 'add_teacher':
                group.teachers.add(user)
            elif action == 'remove_teacher':
                group.teachers.remove(user)
            else:
                return Response({"error": "Invalid action."}, status=status.HTTP_400_BAD_REQUEST)
        except Exception as e:
            return Response({"error": f"Could not perform action: {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        serializer = GroupSerializer(group)
        return Response(serializer.data, status=status.HTTP_200_OK)


class ProjectProgressLogListView(generics.ListAPIView):
    """
    Returns a list of all progress updates for a single project.
    (For teachers to review)
    """
    permission_classes = [IsAuthenticated, IsProjectMemberOrTeacher]
    serializer_class = ProgressUpdateSerializer

    def get_queryset(self):
        project_id = self.kwargs.get('project_id')
        return ProgressUpdate.objects.filter(project_id=project_id)

    def list(self, request, *args, **kwargs):
        project_id = self.kwargs.get('project_id')
        project = get_object_or_404(Project, id=project_id)
        for permission in self.get_permissions():
            if not permission.has_object_permission(request, self, project):
                self.permission_denied(request)
        return super().list(request, *args, **kwargs)
# In authentication/views.py

class RequestPasswordResetView(APIView):
    """
    Takes an email, generates an OTP, saves it, and sends it via email.
    """
    permission_classes = [AllowAny] # Anyone can request a reset

    def post(self, request):
        serializer = PasswordResetRequestSerializer(data=request.data)
        if not serializer.is_valid():
            return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

        email = serializer.validated_data['email']
        try:
            user = User.objects.get(email=email)
        except User.DoesNotExist:
            # For security, don't reveal that the user doesn't exist
            return Response({"detail": "If an account exists with this email, an OTP has been sent."}, status=status.HTTP_200_OK)

        # Generate 6-digit OTP
        otp = f"{random.randint(100000, 999999)}"

        # Save to DB (clear old OTPs first)
        PasswordResetOTP.objects.filter(user=user).delete()
        PasswordResetOTP.objects.create(user=user, otp=otp)

        # Send Email
        send_mail(
            subject="Password Reset OTP",
            message=f"Your OTP for Collage project managment system password reset is: {otp}. It expires in 10 minutes.",
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],
            fail_silently=False,
        )

        return Response({"detail": "OTP sent to email."}, status=status.HTTP_200_OK)


class ResetPasswordView(APIView):
    """
    Takes email, OTP, and new password. If OTP matches, resets the password.
    """
    permission_classes = [AllowAny]

    def post(self, request):
        serializer = PasswordResetConfirmSerializer(data=request.data)
        if not serializer.is_valid():
            return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

        email = serializer.validated_data['email']
        otp = serializer.validated_data['otp']
        new_password = serializer.validated_data['new_password']

        try:
            user = User.objects.get(email=email)
        except User.DoesNotExist:
             return Response({"error": "Invalid request."}, status=status.HTTP_400_BAD_REQUEST)

        # Find the OTP
        try:
            otp_record = PasswordResetOTP.objects.get(user=user, otp=otp)
        except PasswordResetOTP.DoesNotExist:
            return Response({"error": "Invalid OTP."}, status=status.HTTP_400_BAD_REQUEST)

        # Check expiration
        if not otp_record.is_valid():
            return Response({"error": "OTP has expired."}, status=status.HTTP_400_BAD_REQUEST)

        # Reset Password
        user.set_password(new_password)
        user.save()

        # Delete the used OTP
        otp_record.delete()

        return Response({"detail": "Password has been reset successfully."}, status=status.HTTP_200_OK)

class ProjectExtractionView(APIView):
    """
    Extracts Title and Abstract from an uploaded project document (PDF/PPT).
    """
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)

    def post(self, request):
        file_obj = request.FILES.get('file')
        if not file_obj:
            return Response({"error": "File is required."}, status=status.HTTP_400_BAD_REQUEST)
        
        extracted_text = ""
        try:
            # 1. Extract Text based on file type
            if file_obj.name.lower().endswith('.pdf'):
                from pypdf import PdfReader
                reader = PdfReader(file_obj)
                for page in reader.pages:
                    extracted_text += page.extract_text() + "\n"
            
            elif file_obj.name.lower().endswith(('.ppt', '.pptx')):
                from pptx import Presentation
                prs = Presentation(file_obj)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            else:
                 return Response({"error": "Unsupported file format. Please upload PDF or PPTX."}, status=status.HTTP_400_BAD_REQUEST)
            
            if not extracted_text.strip():
                 return Response({"error": "Could not extract text from file."}, status=status.HTTP_400_BAD_REQUEST)

            # 2. Call AI Microservice to Parse Title/Abstract
            ai_payload = {"text": extracted_text}
            response = requests.post('http://127.0.0.1:8001/parse-project-text', json=ai_payload)
            
            if response.status_code == 200:
                return Response(response.json(), status=status.HTTP_200_OK)
            else:
                return Response({"error": "AI failed to parse text."}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        except Exception as e:
            return Response({"error": f"Extraction failed: {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# In authentication/views.py

class ProjectArtifactUploadView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)

    def post(self, request, project_id, *args, **kwargs):
        project = get_object_or_404(Project, id=project_id)
        
        # Security check
        if project.submission.student != request.user:
             return Response({"error": "Unauthorized."}, status=status.HTTP_403_FORBIDDEN)

        file_obj = request.FILES.get('image_file')
        description = request.data.get('description', '')

        if not file_obj:
             return Response({"error": "Image file is required."}, status=status.HTTP_400_BAD_REQUEST)

        # 1. Save the artifact first (to get a file path)
        artifact = ProjectArtifact.objects.create(
            project=project,
            image_file=file_obj,
            description=description
        )

        # 2. Send to AI
        try:
            # We use the file path from the saved instance
            analysis = analyzer.analyze_image_artifact(artifact.image_file.path)
            
            artifact.extracted_text = analysis['extracted_text']
            artifact.ai_tags = analysis['tags']
            artifact.save()
            
        except Exception as e:
            print(f"AI Artifact Error: {e}")
            # We still return success, just without AI data
            
        serializer = ProjectArtifactSerializer(artifact)
        return Response(serializer.data, status=status.HTTP_201_CREATED)

class ProjectArtifactListView(generics.ListAPIView):
    permission_classes = [IsAuthenticated]
    serializer_class = ProjectArtifactSerializer

    def get_queryset(self):
        project_id = self.kwargs.get('project_id')
        return ProjectArtifact.objects.filter(project_id=project_id).order_by('-uploaded_at')
class ProjectReportUploadView(APIView):
    """
    Allows a student to upload their final report PDF.
    """
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)

    def post(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        if project.submission.student != request.user:
             return Response({"error": "Unauthorized."}, status=status.HTTP_403_FORBIDDEN)
        
        file_obj = request.FILES.get('final_report')
        if not file_obj:
             return Response({"error": "File is required."}, status=status.HTTP_400_BAD_REQUEST)

        project.final_report = file_obj
        project.save()
        
        # Create a ProgressUpdate to notify teacher and allow approval
        # Find the likely checkpoint for this
        checkpoint = Checkpoint.objects.filter(project=project, title__icontains='documentation').first()
        if not checkpoint:
             checkpoint = Checkpoint.objects.filter(project=project, title__icontains='final report').first()

        ProgressUpdate.objects.create(
            project=project,
            author=request.user,
            update_text=f"Final Report Uploaded: {file_obj.name}", 
            ai_suggested_percentage=100,
            sentiment="Neutral",
            checkpoint=checkpoint,
            status='Pending',
            ai_analysis_result={"feedback": "Final Report uploaded. AI verification skipped for large document. Please review the final project documentation manually."}
        )
        
        return Response({"detail": "Report uploaded successfully. Teacher notified."}, status=status.HTTP_200_OK)


class ProjectReportGradeView(APIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]

    def post(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        
        if not project.final_report:
             return Response({"error": "No report uploaded yet."}, status=status.HTTP_400_BAD_REQUEST)
             
        try:
            # Call the analyzer (returns a dict now)
            result = analyzer.grade_final_report(project.final_report.path)
            
            # Save BOTH fields
            project.ai_report_feedback = result['feedback']
            project.final_report_content = result['content'] # <-- New Field
            project.save()
            
            return Response({"feedback": result['feedback']}, status=status.HTTP_200_OK)

        except google_exceptions.ResourceExhausted:
            return Response(
                {"detail": "AI is busy (Rate Limit). Please wait 1 minute."},
                status=status.HTTP_429_TOO_MANY_REQUESTS
            )
        except Exception as e:
            return Response(
                {"detail": f"Analysis failed: {str(e)}"},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
class ProjectTaskManagerView(APIView):
    """
    GET: Fetch all tasks for a project.
    POST: Auto-generate tasks using AI (if none exist).
    """
    permission_classes = [IsAuthenticated]

    def get(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        # Check permissions (Student owner or Teacher)
        if project.submission.student != request.user and not request.user.role in ['Teacher', 'HOD/Admin']:
             return Response({"error": "Unauthorized"}, status=status.HTTP_403_FORBIDDEN)
             
        tasks = project.tasks.all()
        serializer = TaskSerializer(tasks, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    def post(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        if project.submission.student != request.user:
             return Response({"error": "Only the student owner can generate tasks."}, status=status.HTTP_403_FORBIDDEN)
        
        # Optional: Check if tasks already exist to prevent overwriting
        if project.tasks.exists() and request.data.get('force') != True:
             return Response({"detail": "Tasks already exist."}, status=status.HTTP_200_OK)

        # Call AI
        task_titles = analyzer.generate_project_tasks(project.title, project.abstract)
        
        # Save to DB
        created_tasks = []
        for title in task_titles:
            task = Task.objects.create(project=project, title=title, status='To Do')
            created_tasks.append(task)
            
        serializer = TaskSerializer(created_tasks, many=True)
        return Response(serializer.data, status=status.HTTP_201_CREATED)

class TaskUpdateView(generics.UpdateAPIView):
    """
    Update a single task's status (e.g., dragging from 'To Do' to 'Done').
    """
    permission_classes = [IsAuthenticated]
    queryset = Task.objects.all()
    serializer_class = TaskSerializer
    lookup_url_kwarg = 'task_id'

    def patch(self, request, *args, **kwargs):
        return super().partial_update(request, *args, **kwargs)

class ProjectResumeView(APIView):
    """
    Generates professional resume bullet points for a project.
    """
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        
        if project.submission.student != request.user:
             return Response({"error": "Unauthorized"}, status=status.HTTP_403_FORBIDDEN)
        
        # 1. Gather context (Tasks help the AI know the specific tech stack)
        tasks = project.tasks.all()
        tasks_text = ", ".join([t.title for t in tasks])
        
        # 2. Call AI
        points = analyzer.generate_resume_points(project.title, project.abstract, tasks_text)
        
        # 3. Save to DB so we don't have to re-generate later
        project.ai_resume_points = points
        project.save()
        
        return Response({"points": points}, status=status.HTTP_200_OK)
class TypingUpdateView(APIView):
    """
    POST: Updates the timestamp saying 'I am typing'.
    GET: Returns a list of users currently typing (active in last 3 seconds).
    """
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        # Update or create the typing record for this user
        TypingStatus.objects.update_or_create(
            project=project, user=request.user,
            defaults={'timestamp': timezone.now()}
        )
        return Response({"status": "updated"}, status=status.HTTP_200_OK)

    def get(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        
        # Find anyone who typed in the last 3 seconds
        threshold = timezone.now() - timedelta(seconds=3)
        
        typing_users = TypingStatus.objects.filter(
            project=project,
            timestamp__gte=threshold
        ).exclude(user=request.user) # Don't show myself
        
        usernames = [t.user.username for t in typing_users]
        return Response({"typing_users": usernames}, status=status.HTTP_200_OK)
class CodeReviewView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        
        # Check if user is part of the project team
        if not project.team.members.filter(id=request.user.id).exists():
            return Response({"error": "You are not a member of this project."}, status=status.HTTP_403_FORBIDDEN)

        file_name = request.data.get('file_name')
        code_content = request.data.get('code_content')

        if not file_name or not code_content:
            return Response({"error": "file_name and code_content are required."}, status=status.HTTP_400_BAD_REQUEST)

        # Call AI Microservice
        ai_response = {}
        try:
            # Prepare context
            project_context = f"Project Title: {project.title}\nDescription: {project.abstract}\nCategory: {project.category}"

            response = requests.post(
                "http://127.0.0.1:8001/review-code",
                json={
                    "code": code_content, 
                    "filename": file_name,
                    "context": project_context
                },
                timeout=30
            )
            if response.status_code == 200:
                ai_response = response.json()
            else:
                print(f"AI Error: {response.text}")
        except Exception as e:
            print(f"AI Service Exception: {e}")

        # Create CodeReview entry
        review = CodeReview.objects.create(
            project=project,
            student=request.user,
            file_name=file_name,
            code_content=code_content,
            security_score=ai_response.get('security_score', 0),
            quality_score=ai_response.get('quality_score', 0),
            security_issues=json.dumps(ai_response.get('security_issues', [])),
            optimization_tips=json.dumps(ai_response.get('optimization_tips', [])),
            ai_feedback=ai_response.get('ai_feedback', "AI Review Unavailable")
        )

        serializer = CodeReviewSerializer(review)
        return Response(serializer.data, status=status.HTTP_201_CREATED)

    def get(self, request, project_id):
        # Get all reviews for this project
        project = get_object_or_404(Project, id=project_id)
        
        # Check permissions (Teacher, HOD/Admin, or Team Member)
        is_teacher = request.user.role in ['Teacher', 'HOD/Admin']
        is_member = project.team.members.filter(id=request.user.id).exists()
        
        if not (is_teacher or is_member):
             return Response({"error": "Permission denied."}, status=status.HTTP_403_FORBIDDEN)
             
        reviews = CodeReview.objects.filter(project=project).order_by('-uploaded_at')
        serializer = CodeReviewSerializer(reviews, many=True)
        return Response(serializer.data)
        return Response(serializer.data)

class CheckpointListView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        # Check permissions
        if not (request.user.role in ['Teacher', 'HOD/Admin'] or project.team.members.filter(id=request.user.id).exists()):
            return Response({"error": "Permission denied."}, status=status.HTTP_403_FORBIDDEN)
            
        checkpoints = Checkpoint.objects.filter(project=project).order_by('deadline')
        serializer = CheckpointSerializer(checkpoints, many=True)
        return Response(serializer.data)

class CheckpointGenerationView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        if not project.team.members.filter(id=request.user.id).exists():
             return Response({"error": "Only team members can generate the roadmap."}, status=status.HTTP_403_FORBIDDEN)

        # Call AI Microservice to generate checkpoints
        try:
            # Prepare data for AI
            ai_payload = {
                "title": project.title,
                "description": project.abstract,
                "category": project.category
            }
            
            # Call AI Microservice
            response = requests.post('http://127.0.0.1:8001/generate-checkpoints', json=ai_payload)
            
            if response.status_code == 200:
                checkpoints_data = response.json().get('checkpoints', [])
                
                # Clear existing checkpoints if any (optional, or append?) -> Clear for now to regenerate
                Checkpoint.objects.filter(project=project).delete()
                
                created_checkpoints = []
                for cp in checkpoints_data:
                    # Calculate deadline based on project timeline (simplified for now)
                    # In a real app, AI could suggest relative deadlines (e.g., "Week 2")
                    checkpoint = Checkpoint.objects.create(
                        project=project,
                        title=cp.get('title'),
                        description=cp.get('description'),
                        # deadline=... (logic to set deadline)
                    )
                    created_checkpoints.append(checkpoint)
                
                serializer = CheckpointSerializer(created_checkpoints, many=True)
                return Response(serializer.data, status=status.HTTP_201_CREATED)
            else:
                return Response({"error": "AI service failed to generate checkpoints."}, status=status.HTTP_503_SERVICE_UNAVAILABLE)
                
        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class CheckpointVerificationView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id, checkpoint_id):
        project = get_object_or_404(Project, id=project_id)
        checkpoint = get_object_or_404(Checkpoint, id=checkpoint_id, project=project)
        
        if not project.team.members.filter(id=request.user.id).exists():
             return Response({"error": "Permission denied."}, status=status.HTTP_403_FORBIDDEN)

        proof_text = request.data.get('proof_text')
        
        if not proof_text:
             return Response({"error": "Proof is required."}, status=status.HTTP_400_BAD_REQUEST)

        # Call AI to verify
        try:
            ai_payload = {
                "checkpoint_title": checkpoint.title,
                "checkpoint_description": checkpoint.description,
                "proof_text": proof_text,
                "project_context": project.abstract
            }
            
            response = requests.post('http://127.0.0.1:8001/verify-checkpoint', json=ai_payload)
            
            if response.status_code == 200:
                verification_result = response.json()
                feedback = verification_result.get('feedback', '')
                
                # Create ProgressUpdate with 'Pending' status
                ProgressUpdate.objects.create(
                    project=project,
                    author=request.user,
                    update_text=proof_text, # Store the proof
                    ai_suggested_percentage=verification_result.get('suggested_progress', 0),
                    sentiment="Neutral",
                    checkpoint=checkpoint,
                    status='Pending',
                    ai_analysis_result=verification_result
                )
                
                return Response({
                    "is_approved": False, 
                    "status": "Pending",
                    "detail": "Proof submitted. Waiting for teacher approval.",
                    "feedback": feedback,
                    "checkpoint_id": checkpoint.id
                })
            else:
                 return Response({"error": "AI verification failed."}, status=status.HTTP_503_SERVICE_UNAVAILABLE)

        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class ProgressUpdateDecisionView(APIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]

    def post(self, request, update_id):
        update = get_object_or_404(ProgressUpdate, id=update_id)
        action = request.data.get('action') # 'approve' or 'reject'
        
        if action == 'approve':
            update.status = 'Approved'
            update.save()
            
            # Complete the checkpoint
            if update.checkpoint:
                update.checkpoint.is_completed = True
                update.checkpoint.date_completed = timezone.now()
                update.checkpoint.save()
            
            # Update Project Progress
            new_progress = update.ai_suggested_percentage
            
            # Use Checkpoint Logic to cap progress
            if update.checkpoint:
                # Rule: Progress is capped at 95% unless the 'Documentation' or 'Final Report' checkpoint is complete.
                title_lower = update.checkpoint.title.lower()
                if "documentation" not in title_lower and "final report" not in title_lower:
                     if new_progress > 95:
                         new_progress = 95
                else:
                    # If this IS documentation/final report, and it's approved, we force 100%
                    if new_progress < 100:
                        new_progress = 100 # Ensure documentation completion finishes the project

            if new_progress > update.project.progress_percentage:
                update.project.progress_percentage = new_progress
                update.project.save()
                
            return Response({"status": "Approved"}, status=status.HTTP_200_OK)
            
        elif action == 'reject':
            update.status = 'Rejected'
            update.save()
            return Response({"status": "Rejected"}, status=status.HTTP_200_OK)
            
        return Response({"error": "Invalid action"}, status=status.HTTP_400_BAD_REQUEST)


class TimedAssignmentCreateView(generics.CreateAPIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]
    serializer_class = TimedAssignmentSerializer

    def perform_create(self, serializer):
        assignment = serializer.save(created_by=self.request.user)
        
        # Send email to all students in assigned groups
        student_emails = []
        for group in assignment.assigned_groups.all():
            for student in group.students.all():
                if student.email:
                    student_emails.append(student.email)
        
        if student_emails:
            try:
                send_mail(
                    subject=f"New Assignment: {assignment.title}",
                    message=f"A new timed assignment '{assignment.title}' has been created.\nDuration: {assignment.duration_minutes} minutes.\nDescription: {assignment.description}\n\nLog in to your dashboard to start.",
                    from_email=settings.DEFAULT_FROM_EMAIL,
                    recipient_list=student_emails,
                    fail_silently=True
                )
            except Exception:
                pass

class TimedAssignmentListView(generics.ListAPIView):
    permission_classes = [IsAuthenticated]
    serializer_class = TimedAssignmentSerializer

    def get_queryset(self):
        user = self.request.user
        if user.role in ['Teacher', 'HOD/Admin']:
            return TimedAssignment.objects.filter(created_by=user).order_by('-start_time')
        elif user.role == 'Student':
            # Get assignments for groups the student belongs to
            return TimedAssignment.objects.filter(assigned_groups__students=user).distinct().order_by('-start_time')
        return TimedAssignment.objects.none()

class AssignmentSubmissionView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = [MultiPartParser, FormParser]

    def post(self, request, assignment_id):
        assignment = get_object_or_404(TimedAssignment, id=assignment_id)
        
        if not assignment.is_active:
            return Response({"error": "Assignment time has expired."}, status=status.HTTP_400_BAD_REQUEST)
            
        # Check if user is in assigned group
        user_groups = request.user.student_groups.all()
        assigned_groups = assignment.assigned_groups.all()
        # Intersection of user groups and assignment groups
        # Since these are QuerySets, we can filter
        common_groups = user_groups.filter(id__in=assigned_groups.values_list('id', flat=True))
        
        if not common_groups.exists():
             return Response({"error": "You are not assigned to this task."}, status=status.HTTP_403_FORBIDDEN)
             
        data = request.data.copy()
        data['assignment'] = assignment.id
        # We need to manually set submitted_by because it's read_only in serializer
        # But wait, serializer.save(submitted_by=request.user) is better
        data['group'] = common_groups.first().id 
        
        serializer = AssignmentSubmissionSerializer(data=data)
        if serializer.is_valid():
            submission = serializer.save(submitted_by=request.user)
            
            # Start AI verification in background thread (non-blocking)
            import threading
            
            def verify_in_background():
                try:
                    # --- 1. Fetch Project Context (Applies to ALL assignment types) ---
                    group = submission.group
                    project_context = ""
                    if group:
                        # Find active project for this group via submission relationship
                        active_projects = Project.objects.filter(
                            submission__group=group,
                            submission__status__in=['Approved', 'In Progress']
                        ).first()
                        
                        if active_projects:
                            project_context = f"Project: {active_projects.title}\nAbstract: {active_projects.abstract}\nCategory: {active_projects.category}"

                    # --- 2. Determine Content to Verify ---
                    content_to_verify = submission.text_content or ""
                    image_data = None
                    
                    # Handle File Content Safely
                    if submission.file:
                        try:
                            file_name = submission.file.name.lower()
                            # Check if it's an image
                            if file_name.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp')):
                                import base64
                                submission.file.seek(0)
                                image_bytes = submission.file.read()
                                image_data = base64.b64encode(image_bytes).decode('utf-8')
                                submission.file.seek(0) # Reset
                                content_to_verify += f"\n\n[Image File Uploaded]: {submission.file.name}"
                            else:
                                # Try to read as text (e.g., code files, txt)
                                file_content = submission.file.read().decode('utf-8')
                                submission.file.seek(0) # Reset pointer
                                if not content_to_verify:
                                    content_to_verify = file_content
                                else:
                                    content_to_verify += f"\n\n[Attached File Content]:\n{file_content}"
                        except UnicodeDecodeError:
                            # Binary file (PDF, etc.) - Don't crash!
                            submission.file.seek(0) # Reset pointer
                            content_to_verify += f"\n\n[Binary File Uploaded]: {submission.file.name}\n(AI Note: Assume this file contains the required diagram/report as described by the filename.)"
                        except Exception as e:
                            content_to_verify += f"\n\n[Error reading file]: {str(e)}"

                    # --- 3. Route to Appropriate AI Endpoint ---
                    
                    if assignment.assignment_type == 'Code' and submission.file and not image_data:
                        # --- Code Review Specific Logic (Text Only) ---
                        ai_payload = {
                            'code': content_to_verify, 
                            'filename': submission.file.name,
                            'context': f"{assignment.description}\n\n{project_context}" if project_context else assignment.description
                        }
                        
                        try:
                            ai_response = requests.post(
                                'http://127.0.0.1:8001/review-code', 
                                json=ai_payload,
                                timeout=45
                            )
                            
                            if ai_response.status_code == 200:
                                result = ai_response.json()
                                security_score = result.get('security_score', 0)
                                quality_score = result.get('quality_score', 0)
                                submission.ai_verified = True
                                submission.ai_score = (security_score + quality_score) / 2
                                submission.ai_feedback = f"Security: {security_score}/10, Quality: {quality_score}/10\n{result.get('ai_feedback', '')}"
                            else:
                                submission.ai_feedback = "Code review service unavailable"
                        except requests.Timeout:
                            submission.ai_feedback = "Code review timed out (45s)"
                        except Exception as e:
                            submission.ai_feedback = f"Code review error: {str(e)}"

                    else:
                        # --- Generic Verification (Diagram, Report, Other) ---
                        # Uses the NEW /verify-assignment endpoint (Supports Images!)
                        
                        ai_payload = {
                            'assignment_type': assignment.assignment_type,
                            'description': assignment.description,
                            'text_content': content_to_verify[:5000], # Limit char count
                            'project_context': project_context,
                            'image_data': image_data # Send base64 image if present
                        }
                        
                        try:
                            ai_response = requests.post(
                                'http://127.0.0.1:8001/verify-assignment', 
                                json=ai_payload,
                                timeout=45 # Increased timeout for image processing
                            )
                            
                            if ai_response.status_code == 200:
                                result = ai_response.json()
                                print(f"DEBUG: AI Response: {result}") # VIEW IN TERMINAL
                                submission.ai_verified = result.get('is_approved', False)
                                submission.ai_score = result.get('score', 0)
                                submission.ai_feedback = result.get('feedback', '')
                                print(f"DEBUG: Saving submission {submission.id} with feedback len={len(submission.ai_feedback)}")
                            else:
                                print(f"DEBUG: AI Error Status {ai_response.status_code}: {ai_response.text}")
                                submission.ai_feedback = "Verification service unavailable"
                        except requests.Timeout:
                            print("DEBUG: AI Timeout")
                            submission.ai_feedback = "Verification timed out"
                        except Exception as e:
                            print(f"DEBUG: AI Request Error: {e}")
                            submission.ai_feedback = f"Verification error: {str(e)}"
                    
                    submission.save()
                    print("DEBUG: Submission saved successfully.")
                    
                except Exception as e:
                    # Catch-all for any background errors
                    print(f"DEBUG: Background Thread Crash: {e}")
                    submission.ai_feedback = f"Background verification failed: {str(e)}"
                    submission.save()
            
            # Start background thread
            thread = threading.Thread(target=verify_in_background)
            thread.daemon = True
            thread.start()
            
            # Return immediately - don't wait for AI
            response_data = serializer.data
            response_data['ai_status'] = 'processing'  # Indicate AI is still running
                
            return Response(response_data, status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)


class TeacherGroupListView(generics.ListAPIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]
    serializer_class = GroupSerializer
    queryset = Group.objects.all()

    def get_queryset(self):
        user = self.request.user
        # Strict Teacher Filter: Only show groups I am assigned to.
        return user.teaching_groups.all()


class AssignmentSubmissionsListView(generics.ListAPIView):
    """Teacher view to see all submissions for a specific assignment"""
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]
    serializer_class = AssignmentSubmissionSerializer
    
    def get_queryset(self):
        assignment_id = self.kwargs.get('assignment_id')
        assignment = get_object_or_404(TimedAssignment, id=assignment_id, created_by=self.request.user)
        return AssignmentSubmission.objects.filter(assignment=assignment).select_related('group', 'submitted_by').order_by('-submitted_at')

# --- TEACHER DASHBOARD REAL DATA VIEWS ---

class TeacherDashboardStatsView(APIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]

    def get(self, request):
        user = request.user
        
        # Base QuerySets
        submissions_qs = ProjectSubmission.objects.all()
        projects_qs = Project.objects.all()
        vivas_qs = VivaSession.objects.all()

        # Strict Teacher Filter
        teacher_groups = user.teaching_groups.all()
        
        # 0. Identify groups that ALREADY have an active or completed project
        groups_with_projects = Project.objects.filter(
            status__in=['In Progress', 'Completed'],
            submission__group__isnull=False
        ).values_list('submission__group_id', flat=True)

        # 1. Pending Approvals: Only show submissions from groups that DON'T have a project yet
        submissions_qs = ProjectSubmission.objects.filter(
            group__in=teacher_groups,
            status='Submitted'
        ).exclude(group_id__in=groups_with_projects)
        
        projects_qs = Project.objects.filter(submission__group__in=teacher_groups)
        vivas_qs = VivaSession.objects.filter(project__submission__group__in=teacher_groups)

        # 1. Pending Approvals (Submissions with status='Submitted')
        pending_approvals = submissions_qs.filter(status='Submitted').count()

        # 2. Active Projects (Projects with status='In Progress')
        active_projects = projects_qs.filter(status='In Progress').count()

        # 3. Active Assignments (Strictly those created by THIS teacher)
        now = timezone.now()
        active_assignments = TimedAssignment.objects.filter(
            created_by=user,
            start_time__lte=now,
            end_time__gte=now
        ).count()

        # 4. Vivas Scheduled (Total VivaSessions for now, or could filter by future date)
        vivas_scheduled = vivas_qs.count()

        # 5. Unappointed Ongoing Projects
        unappointed_ongoing = Project.objects.filter(
            submission__group__teachers__isnull=True,
            status='In Progress'
        ).count()

        return Response({
            "pending_approvals": pending_approvals,
            "active_projects": active_projects,
            "active_assignments": active_assignments,
            "vivas_scheduled": vivas_scheduled,
            "unappointed_ongoing": unappointed_ongoing,
        })

class UnappointedOngoingProjectsView(generics.ListAPIView):
    """
    Returns Ongoing projects that have NO teacher assigned.
    Returns only title, abstract (description), and progress.
    """
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]
    serializer_class = ApprovedProjectSerializer

    def get_queryset(self):
        user = self.request.user
        # Show all 'In Progress' projects EXCEPT those assigned to this teacher
        # This includes:
        # 1. Projects with NO teachers (Unassigned)
        # 2. Projects assigned to OTHER teachers
        return Project.objects.filter(
            status='In Progress'
        ).exclude(
            submission__group__in=user.teaching_groups.all()
        ).order_by('-submission__submitted_at')

class TeacherActivityFeedView(APIView):
    permission_classes = [IsAuthenticated, IsTeacherOrAdmin]

    def get(self, request):
        activities = []
        user = request.user

        # Strict Teacher Filter
        teacher_groups = user.teaching_groups.all()
        submissions_qs = ProjectSubmission.objects.filter(group__in=teacher_groups)
        messages_qs = Message.objects.filter(project__submission__group__in=teacher_groups)
        logs_qs = StudentActivityLog.objects.filter(project__submission__group__in=teacher_groups)

        # 1. Recent Submissions (last 5)
        recent_subs = submissions_qs.order_by('-submitted_at')[:5]
        for sub in recent_subs:
            activities.append({
                "id": f"sub_{sub.id}",
                "type": "submission",
                "text": f"New project '{sub.title}' submitted by {sub.student.username}",
                "time": sub.submitted_at, # datetime object
                "timestamp": sub.submitted_at.timestamp()
            })

        # 2. Recent Messages (last 5) - assuming Message model has created_at
        recent_msgs = messages_qs.order_by('-timestamp')[:5]
        for msg in recent_msgs:
            activities.append({
                "id": f"msg_{msg.id}",
                "type": "message",
                "text": f"New message from {msg.sender.username} in {msg.project.title}",
                "time": msg.timestamp,
                "timestamp": msg.timestamp.timestamp()
            })

        # 3. Student Activity Logs (New)
        recent_logs = logs_qs.order_by('-timestamp')[:10]
        for log in recent_logs:
            activities.append({
                "id": f"log_{log.id}",
                "type": "system", # Use 'system' icon for now
                "text": f"{log.student.username}: {log.action}",
                "time": log.timestamp,
                "timestamp": log.timestamp.timestamp()
            })
            
        # Sort by timestamp descending
        activities.sort(key=lambda x: x['timestamp'], reverse=True)
        
        # Take top 20 combined
        final_activities = activities[:20]
        
        # Format time for frontend (e.g., "2 hours ago" logic can be done in frontend, here send ISO)
        for act in final_activities:
            # act['time'] is a datetime, let's just send isoformat
            if hasattr(act['time'], 'isoformat'):
                act['time'] = act['time'].isoformat()
            del act['timestamp']


        return Response(final_activities)

class StudentActivityFeedView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request):
        activities = []
        user = request.user

        # 1. Recent Submissions by the student
        recent_subs = ProjectSubmission.objects.filter(student=user).order_by('-submitted_at')[:5]
        for sub in recent_subs:
            activities.append({
                "id": f"sub_{sub.id}",
                "type": "submission",
                "text": f"You submitted project '{sub.title}'",
                "time": sub.submitted_at,
                "timestamp": sub.submitted_at.timestamp()
            })

        # 2. Recent Messages (received)
        recent_msgs = Message.objects.filter(recipient=user).order_by('-timestamp')[:5]
        for msg in recent_msgs:
            activities.append({
                "id": f"msg_{msg.id}",
                "type": "message",
                "text": f"Message from {msg.sender.username}: {msg.content[:30]}...",
                "time": msg.timestamp,
                "timestamp": msg.timestamp.timestamp()
            })

        # 3. Recent Assignments (assigned to student's group)
        student_groups = user.student_groups.all()
        recent_assignments = TimedAssignment.objects.filter(assigned_groups__in=student_groups).order_by('-start_time')[:5]
        for assign in recent_assignments:
            activities.append({
                "id": f"assign_{assign.id}",
                "type": "assignment",
                "text": f"New assignment: {assign.title}",
                "time": assign.start_time,
                "timestamp": assign.start_time.timestamp()
            })

        # Sort by timestamp descending
        activities.sort(key=lambda x: x['timestamp'], reverse=True)
        
        # Take top 10 combined
        final_activities = activities[:10]
        
        # Format time for frontend
        for act in final_activities:
            if hasattr(act['time'], 'isoformat'):
                act['time'] = act['time'].isoformat()
            del act['timestamp']


        return Response(final_activities)

class ProjectUpdateView(APIView):
    permission_classes = [IsAuthenticated]

    def patch(self, request, project_id, *args, **kwargs):
        project = get_object_or_404(Project, id=project_id)
        
        # Check permissions: User must be the student who submitted or a group member
        if project.submission.student != request.user and (not project.submission.group or request.user not in project.submission.group.students.all()):
             return Response({"error": "You do not have permission to update this project."}, status=status.HTTP_403_FORBIDDEN)

        # Allow updating github_repo_link
        github_link = request.data.get('github_repo_link')
        if github_link is not None:
            project.github_repo_link = github_link
            project.save()
            return Response({'status': 'success', 'github_repo_link': project.github_repo_link})

        return Response({'status': 'no changes made'})

class ProjectAuditView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id):
        # 1. Get Project
        project = get_object_or_404(Project, id=project_id)
        
        # 2. Check Permissions (Student owner, Team member, or Teacher)
        is_member = False
        if project.submission.student == request.user:
            is_member = True
        elif hasattr(project, 'team') and request.user in project.team.members.all():
            is_member = True
        
        if not is_member and request.user.role != 'Teacher':
             return Response({"error": "Not authorized to audit this project"}, status=403)

        # 3. Check for GitHub Link
        if not project.github_repo_link:
            return Response({"error": "No GitHub link provided for this project"}, status=400)

        # 4. Prepare Context
        project_context = f"""
        Title: {project.title}
        Abstract: {project.abstract}
        Category: {project.category}
        Tech Stack: {project.submission.tags}
        """

        # 5. Call AI Service
        ai_payload = {
            "github_repo_link": project.github_repo_link,
            "project_context": project_context
        }
        
        try:
            # Call the new /audit-code endpoint
            response = requests.post("http://127.0.0.1:8001/audit-code", json=ai_payload)
            
            if response.status_code == 200:
                result = response.json()
                
                # 6. Save Valid Results
                if "error" not in result:
                    project.audit_security_score = result.get('security_score', 0)
                    project.audit_quality_score = result.get('quality_score', 0)
                    project.audit_report = result # Save full JSON
                    project.last_audit_date = timezone.now()
                    project.save()
                    
                    return Response(result)
                else:
                    return Response(result, status=500)
            else:
                return Response({"error": "AI Service unavailable"}, status=503)

        except Exception as e:
            return Response({"error": str(e)}, status=500)

class ProjectDocsView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        
        # Check permissions
        if not (request.user == project.submission.student or project.team.members.filter(id=request.user.id).exists() or request.user.role == 'Teacher'):
             return Response({"error": "Unauthorized"}, status=403)

        if not project.github_repo_link:
            return Response({"error": "No GitHub link provided"}, status=400)

        ai_payload = {
            "github_repo_link": project.github_repo_link,
            "project_context": f"Title: {project.title}\nAbstract: {project.abstract}"
        }
        
        try:
            response = requests.post("http://127.0.0.1:8001/generate-docs", json=ai_payload)
            if response.status_code == 200:
                return Response(response.json())
            else:
                return Response({"error": "AI Service unavailable"}, status=503)
        except Exception as e:
            return Response({"error": str(e)}, status=500)

class ProjectIssuesView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        
        # Check permissions
        if not (request.user == project.submission.student or project.team.members.filter(id=request.user.id).exists() or request.user.role == 'Teacher'):
             return Response({"error": "Unauthorized"}, status=403)
             
        if not project.github_repo_link:
            return Response({"error": "No GitHub link provided"}, status=400)

        fake_context = f"Title: {project.title}" # Issues don't strictly need context but endpoint expects AuditCodeIn
        
        ai_payload = {
            "github_repo_link": project.github_repo_link,
            "project_context": fake_context
        }
        
        try:
            response = requests.post("http://127.0.0.1:8001/analyze-issues", json=ai_payload)
            if response.status_code == 200:
                return Response(response.json())
            else:
                return Response({"error": "AI Service unavailable"}, status=503)

        except Exception as e:
            return Response({"error": str(e)}, status=500)

class ProjectAutoFixView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id):
        # 1. Get Project & Check Permissions
        project = get_object_or_404(Project, id=project_id)
        if project.submission.student != request.user:
             # Also allow team members if implemented
             if not (hasattr(project, 'team') and request.user in project.team.members.all()):
                  return Response({"error": "Not authorized"}, status=403)

        # 2. Proxy to AI Microservice
        if not project.github_repo_link:
            return Response({"error": "No GitHub link provided"}, status=400)

        # Payload from frontend: issue_title, issue_description, file_path
        ai_payload = {
            "github_repo_link": project.github_repo_link,
            "issue_title": request.data.get("issue_title"),
            "issue_description": request.data.get("issue_description"),
            "file_path": request.data.get("file_path"),
            "project_context": f"Title: {project.title}\nAbstract: {project.abstract}"
        }
        
        try:
            response = requests.post("http://127.0.0.1:8001/auto-fix", json=ai_payload)
            if response.status_code == 200:
                return Response(response.json())
            else:
                return Response(response.json(), status=response.status_code)
        except Exception as e:
            return Response({"error": str(e)}, status=500)

class TeamMemberView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id, *args, **kwargs):
        """Add a member to the team"""
        try:
            project = Project.objects.get(id=project_id)
            # Permission check: Only Leader or Teacher/Admin can add members
            is_leader = (request.user == project.submission.student)
            is_teacher_admin = (request.user.role in ['Teacher', 'HOD/Admin'] or request.user.is_superuser)

            if not (is_leader or is_teacher_admin):
                return Response({"detail": "Only the Project Leader or Teachers can add members."}, status=status.HTTP_403_FORBIDDEN)
            
            username = request.data.get('username')
            if not username:
                return Response({"detail": "Username is required."}, status=status.HTTP_400_BAD_REQUEST)
            
            try:
                new_member = User.objects.get(username=username)
            except User.DoesNotExist:
                return Response({"detail": "User not found."}, status=status.HTTP_404_NOT_FOUND)
            
            if new_member.role != 'Student':
                return Response({"detail": "Only students can be added to the team."}, status=status.HTTP_400_BAD_REQUEST)
            
            if new_member in project.team.members.all():
                return Response({"detail": "User is already in the team."}, status=status.HTTP_400_BAD_REQUEST)
            
            # OPTIONAL: Check if student is already in ANOTHER active team?
            # For now, we'll allow multiple teams or assume frontend checks, but let's be safe:
            # if new_member.active_projects.exclude(status='Archived').exists():
            #    return Response({"detail": "User is already in an active project."}, status=status.HTTP_400_BAD_REQUEST)

            project.team.members.add(new_member)
            return Response({"detail": f"{username} added to the team successfully."}, status=status.HTTP_200_OK)

        except Project.DoesNotExist:
            return Response({"detail": "Project not found."}, status=status.HTTP_404_NOT_FOUND)

    def delete(self, request, project_id, *args, **kwargs):
        """Remove a member from the team"""
        try:
            project = Project.objects.get(id=project_id)
            
            # Permission check: Only Leader or Teacher/Admin can remove members
            is_leader = (request.user == project.submission.student)
            is_teacher_admin = (request.user.role in ['Teacher', 'HOD/Admin'] or request.user.is_superuser)

            if not (is_leader or is_teacher_admin):
                return Response({"detail": "Only the Project Leader or Teachers can remove members."}, status=status.HTTP_403_FORBIDDEN)
            
            user_id = request.data.get('user_id')
            # Look for user_id in query params if not in body (for DELETE requests)
            if not user_id:
                user_id = request.query_params.get('user_id')

            if not user_id:
                return Response({"detail": "User ID is required."}, status=status.HTTP_400_BAD_REQUEST)
            
            try:
                member_to_remove = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({"detail": "User not found."}, status=status.HTTP_404_NOT_FOUND)

            if member_to_remove == project.submission.student:
                return Response({"detail": "Cannot remove the project leader (submitter)."}, status=status.HTTP_400_BAD_REQUEST)

            if member_to_remove not in project.team.members.all():
                return Response({"detail": "User is not in the team."}, status=status.HTTP_400_BAD_REQUEST)

            project.team.members.remove(member_to_remove)
            return Response({"detail": f"{member_to_remove.username} removed from the team."}, status=status.HTTP_200_OK)

        except Project.DoesNotExist:
            return Response({"detail": "Project not found."}, status=status.HTTP_404_NOT_FOUND)

class StudentMyProjectView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request, *args, **kwargs):
        user = request.user
        # 1. Try to find a Project where the user is a team member
        project = Project.objects.filter(team__members=user).first()
        if project:
            return Response({
                "id": project.id,
                "title": project.title,
                "abstract_text": project.abstract,
                "github_link": project.github_repo_link,
                "status": project.status
            }, status=status.HTTP_200_OK)

        # 2. Fallback: Find the latest ProjectSubmission by the user
        submission = ProjectSubmission.objects.filter(student=user).order_by('-submitted_at').first()
        if submission:
            # Check if this submission has a related project (even if user isn't in team explicitly yet)
            if hasattr(submission, 'project'):
                 return Response({
                    "id": submission.project.id,
                    "title": submission.project.title,
                    "abstract_text": submission.project.abstract,
                    "github_link": submission.project.github_repo_link,
                    "status": submission.project.status
                }, status=status.HTTP_200_OK)
            
            return Response({
                "id": None, # No project ID yet
                "submission_id": submission.id,
                "title": submission.title,
                "abstract_text": submission.abstract_text,
                "github_link": None,
                "status": submission.status
            }, status=status.HTTP_200_OK)

        return Response({"detail": "No active project found."}, status=status.HTTP_404_NOT_FOUND)

class StudentActivityLogView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        student = request.user
        action = request.data.get('action')
        project_id = request.data.get('project_id')
        details = request.data.get('details', {})

        if not action:
             return Response({"error": "Action is required."}, status=status.HTTP_400_BAD_REQUEST)

        print(f"DEBUG: Logging Activity: {student.username} - {action}") # Terminal Log

        project = None
        if project_id:
            project = get_object_or_404(Project, id=project_id)
        
        log = StudentActivityLog.objects.create(
            student=student,
            project=project,
            action=action,
            details=details
        )
        

        return Response(StudentActivityLogSerializer(log).data, status=status.HTTP_201_CREATED)

# --- RAG Chat ---
class ProjectChatCodebaseView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, project_id):
        project = get_object_or_404(Project, id=project_id)
        
        # Check if user has access (Leader or Team)
        is_member = project.submission.student == request.user
        if hasattr(project, 'team') and request.user in project.team.members.all():
            is_member = True
            
        if not is_member:
             return Response({"error": "Not authorized"}, status=403)

        if not project.github_repo_link:
            return Response({"error": "No GitHub link provided"}, status=400)

        ai_payload = {
            "github_repo_link": project.github_repo_link,
            "query": request.data.get("query"),
            "project_context": f"Title: {project.title}"
        }
        
        try:
            response = requests.post("http://127.0.0.1:8001/chat-codebase", json=ai_payload)
            if response.status_code == 200:
                return Response(response.json())
            else:
                return Response(response.json(), status=response.status_code)
        except Exception as e:
            return Response({"error": str(e)}, status=500)
